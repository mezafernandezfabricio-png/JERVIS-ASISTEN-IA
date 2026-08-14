# -*- coding: utf-8 -*-
"""
create_document.py — Creador profesional de documentos Word (.docx), Texto (.txt, .md),
Presentaciones (.pptx), Tablas (.csv, .json) y Páginas Web (.html) para JARVIS.
Crea documentos completos, extensos, estructurados y de alta calidad.
Guarda SIEMPRE en el Escritorio del usuario con apertura automática.
"""

import os
import sys
import re
import json
from pathlib import Path
from datetime import datetime

BASE_DIR = Path(__file__).resolve().parent.parent
API_FILE = BASE_DIR / "config" / "api_keys.json"

def _get_desktop_dir() -> Path:
    desktop = Path.home() / "Desktop"
    if not desktop.exists():
        onedrive_desktop = Path.home() / "OneDrive" / "Desktop"
        if onedrive_desktop.exists():
            return onedrive_desktop
        desktop.mkdir(parents=True, exist_ok=True)
    return desktop

def _get_api_key() -> str:
    try:
        if API_FILE.exists():
            return json.loads(API_FILE.read_text(encoding="utf-8")).get("gemini_api_key", "")
    except Exception:
        pass
    return ""

def _sanitize_filename(name: str) -> str:
    cleaned = re.sub(r'[\\/*?:"<>|]', "", name).strip()
    return cleaned if cleaned else f"Documento_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

def _enrich_content_if_needed(title: str, content: str, subtitle: str, doc_format: str) -> str:
    """Si el contenido es corto o es solo un tema, genera un documento extenso y completo con IA."""
    words = content.strip().split()
    if len(words) >= 40 and "\n" in content:
        return content

    api_key = _get_api_key()
    if not api_key:
        return content

    try:
        from google import genai
        client = genai.Client(api_key=api_key)
        prompt = (
            f"Actúa como un redactor profesional de JARVIS. Genera el contenido COMPLETO, EXTENSO, ESTRUCTURADO y DETALLADO para un documento en formato '{doc_format}'.\n"
            f"Título del documento: '{title}'\n"
            f"Tema / Indicación del usuario: '{content}'\n"
            f"Subtítulo / Contexto: '{subtitle}'\n\n"
            "REQUISITOS OBLIGATORIOS:\n"
            "1. Redacta en español formal y claro.\n"
            "2. Estructura el documento usando títulos '# ', subtítulos '## ' y viñetas '- '.\n"
            "3. Incluye: Introducción, Desarrollo exhaustivo con secciones temáticas, ejemplos prácticos, puntos clave, y Conclusiones.\n"
            "4. Escribe contenido 100% real, verídico y profesional sin usar resúmenes superficiales ni placeholders."
        )
        resp = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        if resp.text and len(resp.text.strip()) > 100:
            return resp.text.strip()
    except Exception:
        pass
    return content

def create_document(parameters: dict = None, player=None, speak=None, **kwargs) -> str:
    """
    Crea un documento en el Escritorio con formato profesional y contenido extenso.
    Parámetros:
        - title / titulo: Título o nombre del archivo
        - content / texto / body: Contenido o tema a desarrollar
        - format / formato / extension: Tipo de archivo ('docx', 'txt', 'md', 'pptx', 'csv', 'json', 'html')
        - subtitle / subtitulo: Subtítulo opcional
        - open_file / abrir: Si debe abrir el archivo tras crearlo (default True)
    """
    params = parameters or {}
    title = params.get("title") or params.get("titulo") or "Documento JARVIS"
    content = params.get("content") or params.get("texto") or params.get("body") or params.get("text") or "Documento generado por JARVIS."
    subtitle = params.get("subtitle") or params.get("subtitulo") or ""
    author = params.get("author") or params.get("autor") or "JARVIS"
    doc_format = str(params.get("format") or params.get("formato") or params.get("extension") or "").lower().strip().replace(".", "")
    open_file = params.get("open_file", params.get("abrir", True))
    
    clean_title = _sanitize_filename(title)
    if not doc_format:
        for ext in ["docx", "doc", "txt", "md", "pptx", "csv", "json", "html"]:
            if clean_title.lower().endswith(f".{ext}"):
                doc_format = ext
                clean_title = clean_title[:-len(ext)-1]
                break
    if not doc_format:
        doc_format = "docx"

    if player:
        try: player.write_log(f"📝 Redactando documento {doc_format.upper()}: '{clean_title}'...")
        except: pass

    # Enriquecer y generar contenido completo
    content = _enrich_content_if_needed(title, content, subtitle, doc_format)

    desktop = _get_desktop_dir()
    file_path = desktop / f"{clean_title}.{doc_format}"

    try:
        if doc_format in ["docx", "doc"]:
            _create_word_doc(file_path, title, subtitle, content, author)
        elif doc_format == "pptx":
            _create_pptx_doc(file_path, title, subtitle, content)
        elif doc_format in ["txt", "md"]:
            _create_text_doc(file_path, title, subtitle, content, author)
        elif doc_format == "html":
            _create_html_doc(file_path, title, subtitle, content, author)
        elif doc_format == "json":
            _create_json_doc(file_path, title, content)
        elif doc_format == "csv":
            _create_csv_doc(file_path, title, content)
        else:
            file_path.write_text(str(content), encoding="utf-8")

        if player:
            try: player.write_log(f"📁 Documento guardado en Escritorio: {file_path.name}")
            except: pass

        if open_file and file_path.exists():
            try: os.startfile(str(file_path))
            except: pass

        return f"Documento creado exitosamente en tu Escritorio: '{file_path.name}' con toda la información solicitada."

    except Exception as e:
        fallback_path = desktop / f"{clean_title}.txt"
        fallback_path.write_text(f"{title}\n\n{content}", encoding="utf-8")
        if open_file and fallback_path.exists():
            try: os.startfile(str(fallback_path))
            except: pass
        return f"Documento guardado en tu Escritorio: '{fallback_path.name}'"

def _create_word_doc(file_path: Path, title: str, subtitle: str, content: str, author: str):
    """Genera documento Word con estilos, títulos y formato limpio."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Título
        h1 = doc.add_heading(title, level=0)
        h1.alignment = WD_ALIGN_PARAGRAPH.LEFT
        
        if subtitle:
            p_sub = doc.add_paragraph()
            r_sub = p_sub.add_run(subtitle)
            r_sub.italic = True
            r_sub.font.color.rgb = RGBColor(100, 116, 139)
            
        p_meta = doc.add_paragraph()
        r_meta = p_meta.add_run(f"Autor: {author}  |  Fecha: {datetime.now().strftime('%d/%m/%Y %H:%M')}")
        r_meta.font.size = Pt(9)
        r_meta.font.color.rgb = RGBColor(148, 163, 184)
        
        doc.add_paragraph("—" * 35)
        
        for line in str(content).split("\n"):
            line_str = line.strip()
            if not line_str:
                continue
            if line_str.startswith("# "):
                doc.add_heading(line_str[2:], level=1)
            elif line_str.startswith("## "):
                doc.add_heading(line_str[3:], level=2)
            elif line_str.startswith("### "):
                doc.add_heading(line_str[4:], level=3)
            elif line_str.startswith("- ") or line_str.startswith("* "):
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.25)
                r_b = p.add_run("• ")
                r_b.bold = True
                _add_formatted_runs(p, line_str[2:])
            elif re.match(r'^\d+\.\s', line_str):
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.25)
                _add_formatted_runs(p, line_str)
            else:
                p = doc.add_paragraph()
                _add_formatted_runs(p, line_str)
                
        doc.save(str(file_path))
    except Exception:
        html_content = f"""<html xmlns:o='urn:schemas-microsoft-com:office:office' xmlns:w='urn:schemas-microsoft-com:office:word' xmlns='http://www.w3.org/TR/REC-html40'>
<head><title>{title}</title><style>body {{ font-family: Calibri, sans-serif; margin: 40px; color: #1e293b; }} h1 {{ color: #0f172a; }}</style></head>
<body><h1>{title}</h1>{f'<h3><i>{subtitle}</i></h3>' if subtitle else ''}<p><small>Autor: {author} | {datetime.now().strftime('%d/%m/%Y')}</small></p><hr/><p>{content.replace(chr(10), '<br>')}</p></body></html>"""
        file_path.write_text(html_content, encoding="utf-8")

def _add_formatted_runs(paragraph, text: str):
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        else:
            paragraph.add_run(part)

def _create_pptx_doc(file_path: Path, title: str, subtitle: str, content: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if subtitle:
            title_slide.placeholders[1].text = subtitle
            
        lines = str(content).split("\n")
        current_slide = None
        tf = None
        
        for line in lines:
            line_str = line.strip()
            if not line_str: continue
            if line_str.startswith("# ") or current_slide is None:
                h_text = line_str[2:] if line_str.startswith("# ") else title
                current_slide = prs.slides.add_slide(prs.slide_layouts[1])
                current_slide.shapes.title.text = h_text
                tf = current_slide.placeholders[1].text_frame
            else:
                p = tf.add_paragraph()
                p.text = line_str
                p.level = 0 if not (line_str.startswith("- ") or line_str.startswith("* ")) else 1
                
        prs.save(str(file_path))
    except Exception:
        file_path.write_text(f"# {title}\n\n{content}", encoding="utf-8")

def _create_text_doc(file_path: Path, title: str, subtitle: str, content: str, author: str):
    header = f"{title}\n" + ("=" * len(title)) + f"\nAutor: {author} | Fecha: {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n"
    if subtitle: header += f"{subtitle}\n" + ("-" * len(subtitle)) + "\n\n"
    file_path.write_text(header + content, encoding="utf-8")

def _create_html_doc(file_path: Path, title: str, subtitle: str, content: str, author: str):
    html = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <title>{title}</title>
    <style>
        body {{ font-family: 'Segoe UI', Arial, sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; color: #1e293b; line-height: 1.6; background: #f8fafc; }}
        .card {{ background: white; border: 1px solid #e2e8f0; border-radius: 12px; padding: 32px; box-shadow: 0 4px 6px -1px rgb(0 0 0 / 0.1); }}
        h1 {{ color: #0f172a; border-bottom: 2px solid #d97706; padding-bottom: 8px; margin-top: 0; }}
        h2 {{ color: #1e293b; margin-top: 24px; }}
        .meta {{ color: #64748b; font-size: 0.9em; margin-bottom: 20px; }}
        ul {{ padding-left: 20px; }}
        li {{ margin-bottom: 6px; }}
    </style>
</head>
<body>
    <div class="card">
        <h1>{title}</h1>
        {f'<h3><i>{subtitle}</i></h3>' if subtitle else ''}
        <div class="meta">Autor: {author} • Fecha: {datetime.now().strftime('%d/%m/%Y %H:%M')}</div>
        <div>{content.replace(chr(10), '<br>')}</div>
    </div>
</body>
</html>"""
    file_path.write_text(html, encoding="utf-8")

def _create_json_doc(file_path: Path, title: str, content: str):
    try:
        data = json.loads(content)
    except:
        data = {"title": title, "date": datetime.now().isoformat(), "content": content}
    file_path.write_text(json.dumps(data, indent=4, ensure_ascii=False), encoding="utf-8")

def _create_csv_doc(file_path: Path, title: str, content: str):
    file_path.write_text(content, encoding="utf-8")

def document_creator(parameters: dict = None, player=None, speak=None, **kwargs) -> str:
    return create_document(parameters, player, speak, **kwargs)

def document_manager(parameters: dict = None, player=None, speak=None, **kwargs) -> str:
    return create_document(parameters, player, speak, **kwargs)
